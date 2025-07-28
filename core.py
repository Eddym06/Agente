"""
Core module - Lógica de negocio y herramientas del agente
Contiene todas las funciones principales para generación de documentos,
web scraping, integración con LLM y otras herramientas avanzadas.
"""

import os
import json
import yaml
import logging
import requests
from datetime import datetime
from typing import Dict, List, Any, Optional
from pathlib import Path

# Imports para generación de documentos
try:
    from docx import Document
    from docx.shared import Inches
    from pptx import Presentation
    from pptx.util import Inches as PptxInches
except ImportError:
    print("Warning: python-docx or python-pptx not installed")

# Imports para web scraping
try:
    from bs4 import BeautifulSoup
except ImportError:
    print("Warning: beautifulsoup4 not installed")

# Import para OpenAI
try:
    import openai
except ImportError:
    print("Warning: openai package not installed")


class AgentCore:
    """
    Clase principal que contiene toda la lógica de negocio del agente.
    Maneja la configuración, logging, y ejecución de todas las herramientas.
    """
    
    def __init__(self, config_path: str = "config.yaml"):
        """
        Inicializa el core del agente con la configuración especificada.
        
        Args:
            config_path: Ruta al archivo de configuración YAML
        """
        self.config = self._load_config(config_path)
        self.logger = self._setup_logging()
        self._create_directories()
        self._setup_llm()
        
    def _load_config(self, config_path: str) -> Dict[str, Any]:
        """Carga la configuración desde el archivo YAML."""
        try:
            with open(config_path, 'r', encoding='utf-8') as file:
                return yaml.safe_load(file)
        except FileNotFoundError:
            print(f"Warning: Config file {config_path} not found, using defaults")
            return self._get_default_config()
        except yaml.YAMLError as e:
            print(f"Error loading config: {e}")
            return self._get_default_config()
    
    def _get_default_config(self) -> Dict[str, Any]:
        """Retorna configuración por defecto si no se puede cargar el archivo."""
        return {
            'app': {'name': 'Agente Personalizado', 'version': '1.0.0'},
            'paths': {
                'documents_output': './documents',
                'presentations_output': './presentations',
                'logs': './logs',
                'temp': './temp'
            },
            'llm': {'provider': 'lm_studio'},
            'logging': {'level': 'INFO', 'max_lines': 1000}
        }
    
    def _setup_logging(self) -> logging.Logger:
        """Configura el sistema de logging."""
        logger = logging.getLogger('AgentCore')
        logger.setLevel(getattr(logging, self.config.get('logging', {}).get('level', 'INFO')))
        
        # Crear directorio de logs si no existe
        os.makedirs(self.config['paths']['logs'], exist_ok=True)
        
        # Handler para archivo
        file_handler = logging.FileHandler(
            os.path.join(self.config['paths']['logs'], 'agent.log'),
            encoding='utf-8'
        )
        file_handler.setLevel(logging.DEBUG)
        
        # Formato del log
        formatter = logging.Formatter('%(asctime)s - %(name)s - %(levelname)s - %(message)s')
        file_handler.setFormatter(formatter)
        
        logger.addHandler(file_handler)
        return logger
    
    def _create_directories(self):
        """Crea los directorios necesarios para la aplicación."""
        for path_key, path_value in self.config['paths'].items():
            os.makedirs(path_value, exist_ok=True)
            self.logger.info(f"Directory created/verified: {path_value}")
    
    def _setup_llm(self):
        """Configura la conexión con el modelo de lenguaje."""
        llm_config = self.config.get('llm', {})
        self.llm_provider = llm_config.get('provider', 'lm_studio')
        
        if self.llm_provider == 'openai':
            api_key = llm_config.get('openai', {}).get('api_key', '')
            if api_key:
                openai.api_key = api_key
                self.logger.info("OpenAI API configured")
            else:
                self.logger.warning("OpenAI API key not provided")
        elif self.llm_provider == 'lm_studio':
            self.lm_studio_url = llm_config.get('lm_studio', {}).get('base_url', 'http://localhost:1234/v1')
            self.logger.info(f"LM Studio configured: {self.lm_studio_url}")

    # === GENERACIÓN DE DOCUMENTOS ===
    
    def generate_word_document(self, title: str, content: str, filename: str = None) -> str:
        """
        Genera un documento de Word con el contenido especificado.
        
        Args:
            title: Título del documento
            content: Contenido principal del documento
            filename: Nombre del archivo (opcional)
            
        Returns:
            Ruta del archivo generado
        """
        try:
            # Crear documento
            doc = Document()
            
            # Agregar título
            title_paragraph = doc.add_heading(title, 0)
            
            # Agregar fecha
            date_paragraph = doc.add_paragraph(f'Generado el: {datetime.now().strftime("%d/%m/%Y %H:%M")}')
            date_paragraph.alignment = 2  # Alineación derecha
            
            # Agregar separador
            doc.add_paragraph('─' * 50)
            
            # Agregar contenido principal
            for line in content.split('\n'):
                if line.strip():
                    doc.add_paragraph(line.strip())
            
            # Generar nombre de archivo si no se proporciona
            if not filename:
                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                filename = f"documento_{timestamp}.docx"
            
            # Asegurar extensión .docx
            if not filename.endswith('.docx'):
                filename += '.docx'
            
            # Ruta completa del archivo
            file_path = os.path.join(self.config['paths']['documents_output'], filename)
            
            # Guardar documento
            doc.save(file_path)
            
            self.logger.info(f"Word document generated: {file_path}")
            return file_path
            
        except Exception as e:
            error_msg = f"Error generating Word document: {str(e)}"
            self.logger.error(error_msg)
            raise Exception(error_msg)
    
    def generate_powerpoint_presentation(self, title: str, slides_content: List[Dict[str, str]], filename: str = None) -> str:
        """
        Genera una presentación de PowerPoint.
        
        Args:
            title: Título de la presentación
            slides_content: Lista de diccionarios con 'title' y 'content' para cada slide
            filename: Nombre del archivo (opcional)
            
        Returns:
            Ruta del archivo generado
        """
        try:
            # Crear presentación
            prs = Presentation()
            
            # Slide de título
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title_placeholder = slide.shapes.title
            subtitle_placeholder = slide.placeholders[1]
            
            title_placeholder.text = title
            subtitle_placeholder.text = f"Generado el {datetime.now().strftime('%d/%m/%Y')}"
            
            # Agregar slides de contenido
            for slide_data in slides_content:
                bullet_slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(bullet_slide_layout)
                
                # Título del slide
                slide.shapes.title.text = slide_data.get('title', 'Sin título')
                
                # Content del slide
                content = slide_data.get('content', '')
                body_shape = slide.shapes.placeholders[1]
                tf = body_shape.text_frame
                
                # Agregar párrafos
                for line in content.split('\n'):
                    if line.strip():
                        p = tf.add_paragraph()
                        p.text = line.strip()
                        p.level = 0
            
            # Generar nombre de archivo si no se proporciona
            if not filename:
                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                filename = f"presentacion_{timestamp}.pptx"
            
            # Asegurar extensión .pptx
            if not filename.endswith('.pptx'):
                filename += '.pptx'
            
            # Ruta completa del archivo
            file_path = os.path.join(self.config['paths']['presentations_output'], filename)
            
            # Guardar presentación
            prs.save(file_path)
            
            self.logger.info(f"PowerPoint presentation generated: {file_path}")
            return file_path
            
        except Exception as e:
            error_msg = f"Error generating PowerPoint presentation: {str(e)}"
            self.logger.error(error_msg)
            raise Exception(error_msg)

    # === WEB SCRAPING ===
    
    def scrape_website(self, url: str, selector: str = None) -> Dict[str, Any]:
        """
        Realiza web scraping de una URL específica.
        
        Args:
            url: URL del sitio web a scrapear
            selector: Selector CSS opcional para extraer contenido específico
            
        Returns:
            Diccionario con los datos extraídos
        """
        try:
            # Headers para simular un navegador real
            headers = {
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
            }
            
            # Realizar petición
            response = requests.get(url, headers=headers, timeout=10)
            response.raise_for_status()
            
            # Parsear HTML
            soup = BeautifulSoup(response.content, 'html.parser')
            
            # Extraer datos básicos
            result = {
                'url': url,
                'status_code': response.status_code,
                'title': soup.title.string if soup.title else 'Sin título',
                'timestamp': datetime.now().isoformat(),
                'content_length': len(response.content)
            }
            
            # Si se especifica un selector, extraer contenido específico
            if selector:
                elements = soup.select(selector)
                result['selected_content'] = [elem.get_text(strip=True) for elem in elements]
            else:
                # Extraer todo el texto visible
                result['full_text'] = soup.get_text(strip=True, separator=' ')
                
                # Extraer enlaces
                links = soup.find_all('a', href=True)
                result['links'] = [{'text': a.get_text(strip=True), 'href': a['href']} for a in links[:20]]  # Limitar a 20 enlaces
                
                # Extraer imágenes
                images = soup.find_all('img', src=True)
                result['images'] = [{'alt': img.get('alt', ''), 'src': img['src']} for img in images[:10]]  # Limitar a 10 imágenes
            
            self.logger.info(f"Website scraped successfully: {url}")
            return result
            
        except requests.RequestException as e:
            error_msg = f"Error making request to {url}: {str(e)}"
            self.logger.error(error_msg)
            raise Exception(error_msg)
        except Exception as e:
            error_msg = f"Error scraping website {url}: {str(e)}"
            self.logger.error(error_msg)
            raise Exception(error_msg)

    # === INTEGRACIÓN CON LLM ===
    
    def query_llm(self, prompt: str, system_message: str = None) -> str:
        """
        Realiza una consulta al modelo de lenguaje configurado.
        
        Args:
            prompt: Pregunta o instrucción para el modelo
            system_message: Mensaje de sistema opcional
            
        Returns:
            Respuesta del modelo de lenguaje
        """
        try:
            if self.llm_provider == 'openai':
                return self._query_openai(prompt, system_message)
            elif self.llm_provider == 'lm_studio':
                return self._query_lm_studio(prompt, system_message)
            else:
                raise Exception(f"Proveedor LLM no soportado: {self.llm_provider}")
                
        except Exception as e:
            error_msg = f"Error querying LLM: {str(e)}"
            self.logger.error(error_msg)
            raise Exception(error_msg)
    
    def _query_openai(self, prompt: str, system_message: str = None) -> str:
        """Consulta a OpenAI API."""
        messages = []
        
        if system_message:
            messages.append({"role": "system", "content": system_message})
        
        messages.append({"role": "user", "content": prompt})
        
        response = openai.ChatCompletion.create(
            model=self.config['llm']['openai']['model'],
            messages=messages,
            max_tokens=1500,
            temperature=0.7
        )
        
        return response.choices[0].message.content
    
    def _query_lm_studio(self, prompt: str, system_message: str = None) -> str:
        """Consulta a LM Studio local."""
        messages = []
        
        if system_message:
            messages.append({"role": "system", "content": system_message})
        
        messages.append({"role": "user", "content": prompt})
        
        payload = {
            "model": self.config['llm']['lm_studio']['model'],
            "messages": messages,
            "max_tokens": 1500,
            "temperature": 0.7
        }
        
        response = requests.post(
            f"{self.lm_studio_url}/chat/completions",
            json=payload,
            timeout=30
        )
        response.raise_for_status()
        
        result = response.json()
        return result['choices'][0]['message']['content']

    # === FUNCIONES DE UTILIDAD ===
    
    def process_text_with_llm(self, text: str, task: str = "analyze") -> str:
        """
        Procesa texto usando el LLM para diferentes tareas.
        
        Args:
            text: Texto a procesar
            task: Tipo de tarea (analyze, summarize, translate, etc.)
            
        Returns:
            Texto procesado
        """
        task_prompts = {
            "analyze": "Analiza el siguiente texto y proporciona un resumen detallado de los puntos clave:",
            "summarize": "Resume el siguiente texto de manera concisa:",
            "translate": "Traduce el siguiente texto al español:",
            "improve": "Mejora la redacción y claridad del siguiente texto:",
            "extract_keywords": "Extrae las palabras clave más importantes del siguiente texto:"
        }
        
        system_message = "Eres un asistente experto en procesamiento de texto. Proporciona respuestas claras y útiles."
        prompt = f"{task_prompts.get(task, task_prompts['analyze'])}\n\n{text}"
        
        return self.query_llm(prompt, system_message)
    
    def get_available_tools(self) -> List[Dict[str, Any]]:
        """
        Retorna la lista de herramientas disponibles desde la configuración.
        
        Returns:
            Lista de herramientas configuradas
        """
        tools = []
        tools_config = self.config.get('tools', {})
        
        for category, category_tools in tools_config.items():
            for tool in category_tools:
                if tool.get('enabled', True):
                    tools.append({
                        'category': category,
                        'name': tool['name'],
                        'description': tool['description']
                    })
        
        return tools
    
    def save_log_to_file(self, content: str, filename: str = None) -> str:
        """
        Guarda contenido en un archivo de log.
        
        Args:
            content: Contenido a guardar
            filename: Nombre del archivo (opcional)
            
        Returns:
            Ruta del archivo guardado
        """
        if not filename:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"log_{timestamp}.txt"
        
        file_path = os.path.join(self.config['paths']['logs'], filename)
        
        with open(file_path, 'w', encoding='utf-8') as file:
            file.write(f"Log generado el: {datetime.now().strftime('%d/%m/%Y %H:%M:%S')}\n")
            file.write("=" * 50 + "\n\n")
            file.write(content)
        
        self.logger.info(f"Log saved to file: {file_path}")
        return file_path


# === FUNCIONES DE EJEMPLO ===

def example_generate_report():
    """Ejemplo de generación de reporte completo."""
    agent = AgentCore()
    
    # Contenido del reporte
    report_content = """
    REPORTE DE ANÁLISIS AUTOMATIZADO
    
    1. INTRODUCCIÓN
    Este reporte ha sido generado automáticamente por el agente personalizado.
    
    2. DATOS RECOPILADOS
    - Fecha de generación: {fecha}
    - Herramientas utilizadas: Generación automática de documentos
    - Estado del sistema: Operativo
    
    3. CONCLUSIONES
    El sistema está funcionando correctamente y todos los módulos están operativos.
    
    4. RECOMENDACIONES
    - Continuar con el monitoreo regular
    - Actualizar configuraciones según sea necesario
    - Mantener copias de seguridad actualizadas
    """.format(fecha=datetime.now().strftime("%d/%m/%Y %H:%M"))
    
    # Generar documento Word
    doc_path = agent.generate_word_document(
        title="Reporte Automatizado",
        content=report_content,
        filename="reporte_ejemplo.docx"
    )
    
    # Generar presentación PowerPoint
    slides = [
        {
            "title": "Introducción",
            "content": "Este es un reporte generado automáticamente\nContiene análisis y recomendaciones"
        },
        {
            "title": "Datos Principales",
            "content": f"Fecha: {datetime.now().strftime('%d/%m/%Y')}\nEstado: Operativo\nHerramientas: Activas"
        },
        {
            "title": "Conclusiones",
            "content": "Sistema funcionando correctamente\nTodos los módulos operativos\nRecomendaciones implementadas"
        }
    ]
    
    ppt_path = agent.generate_powerpoint_presentation(
        title="Reporte Automatizado - Presentación",
        slides_content=slides,
        filename="presentacion_ejemplo.pptx"
    )
    
    return doc_path, ppt_path


if __name__ == "__main__":
    # Ejemplo de uso
    print("Testing AgentCore...")
    
    try:
        agent = AgentCore()
        print(f"Agent initialized successfully")
        print(f"Available tools: {len(agent.get_available_tools())}")
        
        # Test document generation
        doc_path, ppt_path = example_generate_report()
        print(f"Documents generated:")
        print(f"  - Word: {doc_path}")
        print(f"  - PowerPoint: {ppt_path}")
        
    except Exception as e:
        print(f"Error: {e}")